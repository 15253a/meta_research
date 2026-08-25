from __future__ import annotations

from io import BytesIO
import zipfile

from docx import Document
from pptx import Presentation
import pytest

from meta_research.owners.common import OwnerConflict
from meta_research.writing_contract import (
    WRITING_PAPER_INTENT_SCHEMA,
    WRITING_PRESENTATION_INTENT_SCHEMA,
    WRITING_REPORT_INTENT_SCHEMA,
    normalize_writing_intent,
    validate_writing_document,
    writing_document_profile,
)
from meta_research.writing_renderer import (
    WritingRenderedArtifact,
    WritingRendererRegistry,
    default_writing_renderer_registry,
)
from meta_research.writing_skill import (
    WritingSkillRequest,
    _writing_skill_instructions,
    _writing_skill_resources,
)


_CITATION = {
    "citation_ref": "citation:source-1",
    "source_version_ref": "asset_version:source-1",
    "locator": "line:1",
    "claim": "rare morphology remains visible",
    "source_quote": "rare morphology remains visible",
}


def _supported_claim() -> str:
    return (
        "<!-- meta-research-claim:supported refs=citation:source-1 -->\n"
        "rare morphology remains visible [[citation:citation:source-1]]"
    )


def _paper() -> str:
    sections = []
    for role, heading in (
        ("abstract", "Abstract"),
        ("framing", "Introduction"),
        ("methods", "Methods"),
        ("results", "Results"),
        ("discussion", "Discussion"),
        ("conclusion", "Conclusion"),
    ):
        sections.extend(
            (
                "<!-- meta-research-paper-section "
                f"role={role} -->\n## {heading}",
                (
                    "<!-- meta-research-claim:uncertainty -->\n"
                    f"**Uncertainty:** {heading} remains bounded by the snapshot."
                ),
            )
        )
    sections[3] = _supported_claim()
    return "# Morphology paper\n\n" + "\n\n".join(sections) + "\n"


def _presentation() -> str:
    return (
        "# Morphology deck\n\n"
        "<!-- meta-research-structure -->\n## Slide 1: Evidence\n\n"
        f"{_supported_claim()}\n\n"
        "<!-- meta-research-structure -->\n## Slide 2: Limits\n\n"
        "<!-- meta-research-claim:uncertainty -->\n"
        "**Uncertainty:** transfer beyond the frozen sample is unknown.\n"
    )


@pytest.mark.parametrize(
    ("document_type", "schema_ref", "profile_ref", "default_format", "extension"),
    (
        ("report", WRITING_REPORT_INTENT_SCHEMA, "report-v1", "markdown", ".md"),
        ("paper", WRITING_PAPER_INTENT_SCHEMA, "paper-v1", "docx", ".docx"),
        (
            "presentation",
            WRITING_PRESENTATION_INTENT_SCHEMA,
            "presentation-v1",
            "pptx",
            ".pptx",
        ),
    ),
)
def test_document_profiles_normalize_exact_type_specific_intents(
    document_type: str,
    schema_ref: str,
    profile_ref: str,
    default_format: str,
    extension: str,
) -> None:
    intent = {
        "schema_ref": schema_ref,
        "title": "Exact title",
        "audience": "Exact audience",
        "purpose": "Exact purpose",
        "instructions": "  preserve meaning  ",
    }

    normalized = normalize_writing_intent(document_type, intent)

    assert normalized == {**intent, "instructions": "preserve meaning"}
    profile = writing_document_profile(document_type)
    assert profile.profile_ref == profile_ref
    assert profile.default_format == default_format
    assert profile.canonical_extension == extension
    assert profile.canonical_media_type == profile.media_type
    wrong_schema = (
        WRITING_PAPER_INTENT_SCHEMA
        if document_type == "report"
        else WRITING_REPORT_INTENT_SCHEMA
    )
    with pytest.raises(OwnerConflict, match="writing_intent_invalid"):
        normalize_writing_intent(
            document_type,
            {**intent, "schema_ref": wrong_schema},
        )


def test_document_type_is_canonical_and_ppt_is_only_a_ui_label() -> None:
    with pytest.raises(OwnerConflict, match="writing_document_type_invalid"):
        writing_document_profile("ppt")
    with pytest.raises(OwnerConflict, match="writing_document_type_invalid"):
        writing_document_profile("docx")


def test_paper_profile_requires_semantic_roles_without_freezing_one_genre() -> None:
    inventory = validate_writing_document("paper", _paper(), (_CITATION,))

    assert inventory["document_type"] == "paper"
    assert inventory["profile_ref"] == "paper-v1"
    assert inventory["section_roles"] == [
        "abstract",
        "framing",
        "methods",
        "results",
        "discussion",
        "conclusion",
    ]

    # Visible titles are localizable/editorial; the stable role, not an exact
    # English IMRaD heading, carries the validation semantics.
    localized = _paper().replace("## Methods", "## 方法").replace(
        "## Results", "## 结果与分析"
    )
    assert validate_writing_document("paper", localized, (_CITATION,))[
        "section_count"
    ] == 6

    malformed = _paper().replace("role=methods", "role=appendix")
    with pytest.raises(OwnerConflict, match="writing_paper_structure_invalid"):
        validate_writing_document("paper", malformed, (_CITATION,))
    extra_report_heading = _paper().replace(
        "<!-- meta-research-claim:uncertainty -->\n"
        "**Uncertainty:** Methods remains bounded by the snapshot.",
        "<!-- meta-research-structure -->\n## Report-only heading\n\n"
        "<!-- meta-research-claim:uncertainty -->\n"
        "**Uncertainty:** Methods remains bounded by the snapshot.",
    )
    with pytest.raises(OwnerConflict, match="writing_paper_structure_invalid"):
        validate_writing_document("paper", extra_report_heading, (_CITATION,))


def test_paper_profile_accepts_review_semantics_in_the_same_open_profile() -> None:
    review = _paper()
    review = review.replace("role=methods", "role=evidence").replace(
        "## Methods", "## Review protocol"
    )
    review = review.replace("role=results", "role=synthesis").replace(
        "## Results", "## Evidence synthesis"
    )

    inventory = validate_writing_document("paper", review, (_CITATION,))

    assert inventory["section_roles"][2:4] == ["evidence", "synthesis"]


def test_presentation_profile_requires_sequential_content_bearing_slides() -> None:
    inventory = validate_writing_document(
        "presentation", _presentation(), (_CITATION,)
    )

    assert inventory["document_type"] == "presentation"
    assert inventory["slide_count"] == 2

    malformed = _presentation().replace("Slide 2:", "Slide 3:")
    with pytest.raises(
        OwnerConflict, match="writing_presentation_structure_invalid"
    ):
        validate_writing_document("presentation", malformed, (_CITATION,))


def test_report_renderer_preserves_the_legacy_markdown_bytes_exactly() -> None:
    markdown = "# Legacy report\n\n" + _supported_claim() + "\n"
    registry = default_writing_renderer_registry()

    artifact = registry.render("report", markdown, (_CITATION,))

    assert registry.default_format("report") == "markdown"
    assert artifact.content == markdown.encode("utf-8")
    assert artifact.content_hash == artifact.sha256
    assert artifact.file_extension == ".md"


def test_paper_renderer_returns_reproducible_real_docx() -> None:
    registry = default_writing_renderer_registry()

    first = registry.render("paper", _paper(), (_CITATION,))
    second = registry.render("paper", _paper(), (_CITATION,))

    assert registry.default_format("paper") == "docx"
    assert first == second
    assert first.file_extension == ".docx"
    assert first.media_type == (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    assert Document(BytesIO(first.content)).paragraphs[0].text == "Morphology paper"
    _assert_normalized_zip(first.content)


def test_presentation_renderer_returns_reproducible_editable_pptx() -> None:
    registry = default_writing_renderer_registry()

    first = registry.render("presentation", _presentation(), (_CITATION,))
    second = registry.render("presentation", _presentation(), (_CITATION,))

    assert registry.default_format("presentation") == "pptx"
    assert first == second
    assert first.file_extension == ".pptx"
    assert first.media_type == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    deck = Presentation(BytesIO(first.content))
    assert len(deck.slides) == 2
    assert deck.slides[0].shapes.title.text == "Evidence"
    assert "asset_version:source-1 (line:1)" in "\n".join(
        shape.text for shape in deck.slides[0].shapes if shape.has_text_frame
    )
    _assert_normalized_zip(first.content)


def test_renderer_registry_is_open_to_additional_adapters() -> None:
    class _PdfAdapter:
        document_type = "paper"
        renderer_ref = "test/pdf/v1"
        output_format = "pdf"
        media_type = "application/pdf"
        file_extension = ".pdf"

        def render(
            self, markdown: str, citations: tuple[dict[str, str], ...]
        ) -> bytes:
            del markdown, citations
            return b"%PDF-test"

    registry = WritingRendererRegistry(
        (*default_writing_renderer_registry().adapters, _PdfAdapter())
    )

    assert registry.formats("paper") == ("docx", "pdf")
    artifact = registry.render(
        "paper", _paper(), (_CITATION,), output_format="pdf"
    )
    assert isinstance(artifact, WritingRenderedArtifact)
    assert artifact.content == b"%PDF-test"


def test_skill_request_binds_document_type_and_profile_without_new_session() -> None:
    fields = WritingSkillRequest.__dataclass_fields__

    assert fields["document_type"].default == "report"
    assert fields["profile_ref"].default == "report-v1"
    paper_instructions = _writing_skill_instructions("paper")
    presentation_instructions = _writing_skill_instructions("presentation")
    assert "Paper profile" in paper_instructions
    assert "Presentation profile" in presentation_instructions
    assert paper_instructions != presentation_instructions
    assert set(_writing_skill_resources("report")) == {
        "SKILL.md",
        "agents/openai.yaml",
    }
    assert _writing_skill_instructions("report") == _writing_skill_resources(
        "report"
    )["SKILL.md"]


def _assert_normalized_zip(content: bytes) -> None:
    with zipfile.ZipFile(BytesIO(content)) as package:
        assert package.testzip() is None
        assert package.namelist() == sorted(package.namelist())
        assert {item.date_time for item in package.infolist()} == {
            (1980, 1, 1, 0, 0, 0)
        }
