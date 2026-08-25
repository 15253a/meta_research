from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timezone
import hashlib
from io import BytesIO
import re
from typing import Iterable, Protocol
import zipfile

from docx import Document
from docx.document import Document as DocxDocument
from docx.shared import Pt as DocxPt
from pptx import Presentation
from pptx.presentation import Presentation as PptxPresentation
from pptx.util import Inches, Pt as PptxPt

from meta_research.owners.common import OwnerConflict, canonical_hash
from meta_research.writing_contract import (
    WRITING_MAX_OUTPUT_BYTES,
    validate_writing_document,
    writing_document_profile,
)


_ANCHOR = re.compile(r"\[\[citation:([^\]]+)\]\]")
_STRUCTURE = re.compile(
    r"\A<!-- meta-research-structure -->\n(#{2,6}) ([^\n]+)\Z"
)
_PAPER_SECTION = re.compile(
    r"\A<!-- meta-research-paper-section role=[a-z-]+ -->\n"
    r"(#{2,6}) ([^\n]+)\Z"
)
_SLIDE_HEADING = re.compile(r"Slide [1-9][0-9]?: (.+)\Z")
_MARKER_LINE = re.compile(r"\A<!-- meta-research-[^\n]+ -->\n")
_FIXED_OFFICE_TIME = datetime(2000, 1, 1, tzinfo=timezone.utc)
_ZIP_TIME = (1980, 1, 1, 0, 0, 0)


@dataclass(frozen=True)
class WritingRenderedArtifact:
    document_type: str
    profile_ref: str
    renderer_ref: str
    output_format: str
    media_type: str
    file_extension: str
    content: bytes
    content_hash: str
    renderer_input_hash: str
    artifact_hash: str

    @property
    def sha256(self) -> str:
        return self.content_hash


class WritingRendererAdapter(Protocol):
    document_type: str
    renderer_ref: str
    output_format: str
    media_type: str
    file_extension: str

    def render(
        self, markdown: str, citations: tuple[dict[str, str], ...]
    ) -> bytes: ...


class WritingRendererRegistry:
    """Open renderer dispatch behind the three canonical document types."""

    def __init__(self, adapters: Iterable[WritingRendererAdapter]) -> None:
        ordered = tuple(adapters)
        by_key: dict[tuple[str, str], WritingRendererAdapter] = {}
        for adapter in ordered:
            profile = writing_document_profile(adapter.document_type)
            key = (profile.document_type, adapter.output_format)
            if (
                key in by_key
                or not adapter.renderer_ref
                or not adapter.output_format
                or not adapter.media_type
                or not adapter.file_extension.startswith(".")
            ):
                raise OwnerConflict("writing_renderer_registry_invalid")
            by_key[key] = adapter
        self._adapters = ordered
        self._by_key = by_key

    @property
    def adapters(self) -> tuple[WritingRendererAdapter, ...]:
        return self._adapters

    def formats(self, document_type: str) -> tuple[str, ...]:
        profile = writing_document_profile(document_type)
        return tuple(
            adapter.output_format
            for adapter in self._adapters
            if adapter.document_type == profile.document_type
        )

    def default_format(self, document_type: str) -> str:
        profile = writing_document_profile(document_type)
        if (document_type, profile.default_format) not in self._by_key:
            raise OwnerConflict("writing_renderer_unavailable")
        return profile.default_format

    def render(
        self,
        document_type: str,
        markdown: str,
        citations: tuple[dict[str, str], ...],
        *,
        output_format: str | None = None,
    ) -> WritingRenderedArtifact:
        profile = writing_document_profile(document_type)
        validate_writing_document(document_type, markdown, citations)
        selected_format = output_format or self.default_format(document_type)
        adapter = self._by_key.get((document_type, selected_format))
        if adapter is None:
            raise OwnerConflict("writing_render_format_unsupported")
        renderer_input_hash = canonical_hash(
            {
                "document_type": document_type,
                "profile_ref": profile.profile_ref,
                "renderer_ref": adapter.renderer_ref,
                "output_format": selected_format,
                "markdown": markdown,
                "citations": list(citations),
            }
        )
        try:
            content = adapter.render(markdown, citations)
        except OwnerConflict:
            raise
        except Exception as error:
            raise OwnerConflict("writing_renderer_failed") from error
        if (
            not isinstance(content, bytes)
            or not content
            or len(content) > WRITING_MAX_OUTPUT_BYTES
        ):
            raise OwnerConflict("writing_renderer_output_invalid")
        content_hash = hashlib.sha256(content).hexdigest()
        artifact_hash = canonical_hash(
            {
                "renderer_input_hash": renderer_input_hash,
                "renderer_ref": adapter.renderer_ref,
                "output_format": selected_format,
                "media_type": adapter.media_type,
                "file_extension": adapter.file_extension,
                "content_hash": content_hash,
            }
        )
        return WritingRenderedArtifact(
            document_type=document_type,
            profile_ref=profile.profile_ref,
            renderer_ref=adapter.renderer_ref,
            output_format=selected_format,
            media_type=adapter.media_type,
            file_extension=adapter.file_extension,
            content=content,
            content_hash=content_hash,
            renderer_input_hash=renderer_input_hash,
            artifact_hash=artifact_hash,
        )


class _MarkdownRenderer:
    document_type = "report"
    renderer_ref = "meta-research/renderer/markdown/v1"
    output_format = "markdown"
    media_type = "text/markdown; charset=utf-8"
    file_extension = ".md"

    def render(
        self, markdown: str, citations: tuple[dict[str, str], ...]
    ) -> bytes:
        del citations
        return markdown.encode("utf-8")


class _DocxPaperRenderer:
    document_type = "paper"
    renderer_ref = "meta-research/renderer/docx-paper/v1"
    output_format = "docx"
    media_type = (
        "application/vnd.openxmlformats-officedocument."
        "wordprocessingml.document"
    )
    file_extension = ".docx"

    def render(
        self, markdown: str, citations: tuple[dict[str, str], ...]
    ) -> bytes:
        document = Document()
        assert isinstance(document, DocxDocument)
        _set_docx_properties(document, _document_title(markdown))
        normal = document.styles["Normal"]
        normal.font.name = "Aptos"
        normal.font.size = DocxPt(10.5)
        for kind, level, value in _semantic_blocks(markdown):
            if kind == "title":
                document.add_heading(value, level=0)
            elif kind == "heading":
                document.add_heading(value, level=level)
            else:
                document.add_paragraph(value)
        document.add_heading("Citation ledger", level=1)
        for citation in citations:
            document.add_paragraph(
                (
                    f"[{citation['citation_ref']}] "
                    f"{citation['source_version_ref']} — "
                    f"{citation['locator']}"
                ),
                style="List Bullet",
            )
        output = BytesIO()
        document.save(output)
        return _normalize_office_package(output.getvalue())


class _PptxPresentationRenderer:
    document_type = "presentation"
    renderer_ref = "meta-research/renderer/pptx-presentation/v1"
    output_format = "pptx"
    media_type = (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.presentation"
    )
    file_extension = ".pptx"

    def render(
        self, markdown: str, citations: tuple[dict[str, str], ...]
    ) -> bytes:
        deck = Presentation()
        assert isinstance(deck, PptxPresentation)
        _set_pptx_properties(deck, _document_title(markdown))
        deck.slide_width = Inches(13.333)
        deck.slide_height = Inches(7.5)
        citations_by_ref = {
            citation["citation_ref"]: citation for citation in citations
        }
        for title, body_blocks, citation_refs in _presentation_slides(markdown):
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            assert slide.shapes.title is not None
            slide.shapes.title.text = title
            title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
            title_run.font.size = PptxPt(28)
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()
            for ordinal, body in enumerate(body_blocks):
                paragraph = (
                    text_frame.paragraphs[0]
                    if ordinal == 0
                    else text_frame.add_paragraph()
                )
                paragraph.text = body
                paragraph.level = 0
                paragraph.font.size = PptxPt(18)
            if citation_refs:
                citation_box = slide.shapes.add_textbox(
                    Inches(0.55), Inches(6.9), Inches(12.2), Inches(0.35)
                )
                citation_frame = citation_box.text_frame
                citation_frame.clear()
                citation_frame.paragraphs[0].text = " · ".join(
                    (
                        f"[{citation_ref}] "
                        f"{citations_by_ref[citation_ref]['source_version_ref']} "
                        f"({citations_by_ref[citation_ref]['locator']})"
                    )
                    for citation_ref in citation_refs
                )
                citation_frame.paragraphs[0].font.size = PptxPt(8)
        output = BytesIO()
        deck.save(output)
        return _normalize_office_package(output.getvalue())


def default_writing_renderer_registry() -> WritingRendererRegistry:
    return WritingRendererRegistry(
        (_MarkdownRenderer(), _DocxPaperRenderer(), _PptxPresentationRenderer())
    )


def _document_title(markdown: str) -> str:
    first = markdown.strip().split("\n", 1)[0]
    return first[2:]


def _semantic_blocks(markdown: str) -> tuple[tuple[str, int, str], ...]:
    result: list[tuple[str, int, str]] = []
    blocks = re.split(r"\n[ \t]*\n", markdown.strip())
    for ordinal, block in enumerate(blocks):
        block = block.strip()
        if ordinal == 0:
            result.append(("title", 0, block[2:]))
            continue
        structure = _STRUCTURE.fullmatch(block)
        if structure is None:
            structure = _PAPER_SECTION.fullmatch(block)
        if structure is not None:
            result.append(
                ("heading", len(structure.group(1)) - 1, structure.group(2))
            )
            continue
        result.append(("body", 0, _visible_claim_text(block)))
    return tuple(result)


def _presentation_slides(
    markdown: str,
) -> tuple[tuple[str, tuple[str, ...], tuple[str, ...]], ...]:
    slides: list[tuple[str, list[str], list[str]]] = []
    blocks = re.split(r"\n[ \t]*\n", markdown.strip())[1:]
    for raw_block in blocks:
        block = raw_block.strip()
        structure = _STRUCTURE.fullmatch(block)
        if structure is not None:
            match = _SLIDE_HEADING.fullmatch(structure.group(2))
            assert match is not None
            slides.append((match.group(1), [], []))
        else:
            assert slides
            slides[-1][1].append(_visible_claim_text(block))
            for citation_ref in _ANCHOR.findall(block):
                if citation_ref not in slides[-1][2]:
                    slides[-1][2].append(citation_ref)
    return tuple(
        (title, tuple(body), tuple(citation_refs))
        for title, body, citation_refs in slides
    )


def _visible_claim_text(block: str) -> str:
    value = _MARKER_LINE.sub("", block, count=1)
    value = _ANCHOR.sub(lambda match: f"[{match.group(1)}]", value)
    return value.replace("**Inference:**", "Inference:").replace(
        "**Uncertainty:**", "Uncertainty:"
    ).replace("**Evidence gap:**", "Evidence gap:")


def _set_docx_properties(document: DocxDocument, title: str) -> None:
    properties = document.core_properties
    properties.title = title
    properties.subject = "Meta-research accepted Writing artifact"
    properties.author = "Meta-research Writing"
    properties.last_modified_by = "Meta-research Writing"
    properties.created = _FIXED_OFFICE_TIME
    properties.modified = _FIXED_OFFICE_TIME
    properties.revision = 1


def _set_pptx_properties(deck: PptxPresentation, title: str) -> None:
    properties = deck.core_properties
    properties.title = title
    properties.subject = "Meta-research accepted Writing artifact"
    properties.author = "Meta-research Writing"
    properties.last_modified_by = "Meta-research Writing"
    properties.created = _FIXED_OFFICE_TIME
    properties.modified = _FIXED_OFFICE_TIME
    properties.revision = 1


def _normalize_office_package(content: bytes) -> bytes:
    source = BytesIO(content)
    destination = BytesIO()
    with zipfile.ZipFile(source, "r") as package, zipfile.ZipFile(
        destination,
        "w",
        compression=zipfile.ZIP_DEFLATED,
        compresslevel=9,
        strict_timestamps=True,
    ) as normalized:
        for name in sorted(package.namelist()):
            info = zipfile.ZipInfo(name, date_time=_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = 0
            info.external_attr = 0
            normalized.writestr(
                info,
                package.read(name),
                compress_type=zipfile.ZIP_DEFLATED,
            )
    return destination.getvalue()
